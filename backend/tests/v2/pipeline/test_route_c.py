"""Pipeline Route C 테스트 — DocumentToMarkdown + MarkdownToStructured"""
import pytest
from unittest.mock import patch, MagicMock
import io
import sys
import types
from pathlib import Path

from app.services.v2.pipeline.base import PipelineContext
from app.services.v2.pipeline.steps.document_to_md import DocumentToMarkdownStep
from app.services.v2.pipeline.steps.md_to_structured import MarkdownToStructuredStep
from app.services.v2.pipeline.engine import execute_route_c


def make_ctx(spec=None):
    ctx = PipelineContext(dataset_id="test-ds", version_no=1, route="C")
    if spec:
        ctx.spec = spec
    return ctx


SAMPLE_MD_ROW = {
    "filename": "policy.pdf",
    "markdown_text": """# Policy\n\nDate: 2024-01-15\nAuthor: John\nDepartment: Legal""",
}


# ── DocumentToMarkdownStep ─────────────────────────────────────────────

def test_document_to_md_markitdown_bytes():
    """bytes 컨텐츠를 markitdown 전략으로 처리"""
    step = DocumentToMarkdownStep()
    ctx = make_ctx({"document_to_md": {"strategy": "markitdown"}})
    data = [{"filename": "test.txt", "content": b"Hello world content"}]
    result = step.run(ctx, data)
    assert len(result) == 1
    assert "markdown_text" in result[0]
    assert result[0]["extraction_strategy"] == "markitdown"


def test_document_to_md_string_content():
    """문자열 컨텐츠 처리"""
    step = DocumentToMarkdownStep()
    ctx = make_ctx({"document_to_md": {"strategy": "markitdown"}})
    data = [{"filename": "note.txt", "content": "Simple text content"}]
    result = step.run(ctx, data)
    assert result[0]["markdown_text"] != "" or "extraction_error" not in result[0]


def test_document_to_md_ocr_no_paddleocr():
    """OCR disabled/unavailable returns gracefully."""
    step = DocumentToMarkdownStep()
    ctx = make_ctx({"document_to_md": {"strategy": "ocr"}})
    data = [{"filename": "scan.png", "content": b"\x89PNG\r\n"}]
    result = step.run(ctx, data)
    assert len(result) == 1
    assert "markdown_text" in result[0]  # 오류 없이 처리됨


def test_document_to_md_ocr_easyocr_reads_text(monkeypatch):
    """EasyOCR provider extracts text when enabled and dependency is available."""
    fake_easyocr = types.SimpleNamespace()

    class FakeReader:
        def __init__(self, langs, gpu=False):
            self.langs = langs
            self.gpu = gpu

        def readtext(self, path, detail=0, paragraph=True):
            return ["供应商风险", "交付准时率 95%"]

    fake_easyocr.Reader = FakeReader
    monkeypatch.setitem(sys.modules, "easyocr", fake_easyocr)

    step = DocumentToMarkdownStep()
    ctx = make_ctx({"document_to_md": {"strategy": "ocr", "provider": "easyocr"}})
    with patch.object(step, "_load_ocr_runtime_config", return_value={
        "provider": "easyocr",
        "enabled": True,
        "langs": ["ch_sim", "en"],
        "device": "cpu",
    }):
        result = step.run(ctx, [{"filename": "scan.png", "content": b"\x89PNG\r\n\x1a\nfake"}])

    assert "供应商风险" in result[0]["markdown_text"]
    assert result[0]["extraction_method"] == "easyocr"


def test_document_to_md_ocr_disabled_sets_method():
    step = DocumentToMarkdownStep()
    ctx = make_ctx({"document_to_md": {"strategy": "ocr", "provider": "easyocr"}})
    with patch.object(step, "_load_ocr_runtime_config", return_value={"provider": "easyocr", "enabled": False}):
        result = step.run(ctx, [{"filename": "scan.png", "content": b"\x89PNG\r\n\x1a\nfake"}])

    assert result[0]["markdown_text"] == ""
    assert result[0]["extraction_method"] == "easyocr_disabled"


def test_document_to_md_vlm_uses_model_when_configured():
    """VLM strategy calls configured vision model for image input."""
    step = DocumentToMarkdownStep()
    ctx = make_ctx({"document_to_md": {"strategy": "vlm"}})
    data = [{"filename": "chart.png", "content": b"\x89PNG\r\n\x1a\nfake"}]

    with patch.object(step, "_load_vlm_model_config", return_value=MagicMock()), \
         patch.object(step, "_call_vlm", return_value="# Extracted Chart") as mock_call:
        result = step.run(ctx, data)

    assert result[0]["markdown_text"] == "# Extracted Chart"
    mock_call.assert_called_once()


def test_document_to_md_vlm_without_model_falls_back():
    """VLM strategy falls back gracefully when no model is configured."""
    step = DocumentToMarkdownStep()
    ctx = make_ctx({"document_to_md": {"strategy": "vlm"}})
    data = [{"filename": "chart.png", "content": b"\x89PNG\r\n\x1a\nfake"}]

    with patch.object(step, "_load_vlm_model_config", return_value=None), \
         patch.object(step, "_convert_markitdown", return_value="fallback text"):
        result = step.run(ctx, data)

    assert result[0]["markdown_text"] == "fallback text"


def test_document_to_md_vlm_pptx_uses_structured_markdown():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Supply Chain Review"
    textbox = slide.shapes.add_textbox(914400, 1371600, 5486400, 914400)
    textbox.text = "Supplier risk and delivery performance"
    buf = io.BytesIO()
    prs.save(buf)

    step = DocumentToMarkdownStep()
    ctx = make_ctx({"document_to_md": {"strategy": "vlm"}})
    result = step.run(ctx, [{"filename": "review.pptx", "content": buf.getvalue()}])

    assert "# Slide 1" in result[0]["markdown_text"]
    assert "Supply Chain Review" in result[0]["markdown_text"]
    assert "Supplier risk" in result[0]["markdown_text"]


def test_document_to_md_vlm_docx_uses_structured_markdown():
    from docx import Document

    doc = Document()
    doc.add_heading("Procurement Policy", level=1)
    doc.add_paragraph("Approved suppliers must pass review.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Rule"
    table.cell(0, 1).text = "Threshold"
    table.cell(1, 0).text = "OTD"
    table.cell(1, 1).text = "95%"
    buf = io.BytesIO()
    doc.save(buf)

    step = DocumentToMarkdownStep()
    ctx = make_ctx({"document_to_md": {"strategy": "vlm"}})
    result = step.run(ctx, [{"filename": "policy.docx", "content": buf.getvalue()}])

    assert "# Procurement Policy" in result[0]["markdown_text"]
    assert "Approved suppliers" in result[0]["markdown_text"]
    assert "| Rule | Threshold |" in result[0]["markdown_text"]


def test_document_to_md_sets_meta():
    """처리 후 ctx.meta에 통계 기록"""
    step = DocumentToMarkdownStep()
    ctx = make_ctx({"document_to_md": {"strategy": "markitdown"}})
    step.run(ctx, [{"filename": "f.txt", "content": b"test"}])
    assert "document_to_md" in ctx.meta
    assert ctx.meta["document_to_md"]["processed"] == 1


def test_document_to_md_error_graceful():
    """변환 오류 시 빈 markdown_text + error 필드 반환"""
    step = DocumentToMarkdownStep()
    ctx = make_ctx({"document_to_md": {"strategy": "markitdown"}})
    # content 없는 row
    data = [{"filename": "mystery.pdf"}]
    result = step.run(ctx, data)
    assert len(result) == 1
    assert "markdown_text" in result[0]


# ── MarkdownToStructuredStep ───────────────────────────────────────────

def test_md_to_structured_no_schema_passthrough():
    """target_schema 없으면 데이터 그대로 통과"""
    step = MarkdownToStructuredStep()
    ctx = make_ctx({"md_to_structured": {}})
    result = step.run(ctx, [SAMPLE_MD_ROW])
    assert result[0] == SAMPLE_MD_ROW


def test_md_to_structured_with_mock_llm():
    """LLM mock으로 구조화 추출 성공 케이스"""
    step = MarkdownToStructuredStep()
    ctx = make_ctx({
        "md_to_structured": {
            "target_schema": {"date": "날짜", "author": "작성자"},
            "model_id": "test-model",
        }
    })
    with patch("app.services.v2.pipeline.steps.md_to_structured.MarkdownToStructuredStep._extract") as mock_extract:
        mock_extract.return_value = {"date": "2024-01-15", "author": "John"}
        result = step.run(ctx, [SAMPLE_MD_ROW])
    assert result[0]["date"] == "2024-01-15"
    assert result[0]["author"] == "John"
    assert result[0]["structured_extraction_ok"] is True


def test_md_to_structured_llm_failure_graceful():
    """LLM 호출 실패 시 오류 필드 추가 (크래시 없음)"""
    step = MarkdownToStructuredStep()
    ctx = make_ctx({
        "md_to_structured": {
            "target_schema": {"date": "날짜"},
        }
    })
    with patch("app.services.v2.pipeline.steps.md_to_structured.MarkdownToStructuredStep._extract",
               side_effect=RuntimeError("LLM timeout")):
        result = step.run(ctx, [SAMPLE_MD_ROW])
    assert result[0]["structured_extraction_ok"] is False
    assert "structured_extraction_error" in result[0]


def test_md_to_structured_auto_extract_accepts_records_object():
    step = MarkdownToStructuredStep()
    payload = '{"records":[{"record_id":"policy:1","row_type":"section","title":"Policy"}]}'
    with patch("app.services.v2.pipeline.steps.md_to_structured._call_with_model", return_value=payload):
        result = step._auto_extract_with_llm([SAMPLE_MD_ROW], MagicMock())

    assert result is not None
    assert result[0]["record_id"] == "policy:1"
    assert result[0]["row_type"] == "section"


def test_md_to_structured_no_markdown_text():
    """markdown_text 없는 row는 그대로 통과"""
    step = MarkdownToStructuredStep()
    ctx = make_ctx({
        "md_to_structured": {"target_schema": {"date": "날짜"}}
    })
    data = [{"filename": "test.pdf"}]  # markdown_text 없음
    result = step.run(ctx, data)
    assert len(result) == 1


# ── Route C 통합 테스트 ───────────────────────────────────────────────

def test_execute_route_c_basic():
    """Route C 전체 체인 실행 (빈 content → markitdown)"""
    ctx = make_ctx({
        "document_to_md": {"strategy": "markitdown"},
        "md_to_structured": {"target_schema": {"date": "날짜"}},
    })
    data = [{"filename": "test.txt", "content": b"Report date: 2024-01-15"}]
    result, ctx2 = execute_route_c(ctx, data)
    assert len(result) == 1
    assert ctx2.rows_out == 1
    assert "markdown_text" in result[0]


def test_execute_route_c_empty():
    """빈 데이터 입력"""
    ctx = make_ctx({"document_to_md": {"strategy": "markitdown"}})
    result, ctx2 = execute_route_c(ctx, [])
    assert result == []
    assert ctx2.rows_out == 0


def test_execute_route_c_preserves_filename():
    """filename 필드가 유지됨"""
    ctx = make_ctx({"document_to_md": {"strategy": "markitdown"}})
    data = [{"filename": "report.pdf", "content": b"content"}]
    result, _ = execute_route_c(ctx, data)
    assert result[0]["filename"] == "report.pdf"


def test_supply_chain_route_c_real_documents_do_not_drop_content():
    """Supply-chain document fixtures should produce usable markdown/fallback text."""
    base = Path(__file__).resolve().parents[4] / "test_data" / "供应链"
    files = [
        "procurement_policy.docx",
        "supply_chain_review.pptx",
        "supply_chain_strategy.md",
        "warehouse_management.pdf",
    ]
    step = DocumentToMarkdownStep()

    for filename in files:
        ctx = make_ctx({"document_to_md": {"strategy": "vlm"}})
        with patch.object(step, "_load_vlm_model_config", return_value=None):
            result = step.run(ctx, [{"filename": filename, "content": (base / filename).read_bytes()}])

        md = result[0]["markdown_text"]
        assert isinstance(md, str)
        assert md.strip()
        assert "extraction_error" not in result[0]


def test_supply_chain_markdown_rule_based_extracts_rules():
    base = Path(__file__).resolve().parents[4] / "test_data" / "供应链"
    md = (base / "supply_chain_strategy.md").read_text(encoding="utf-8")
    step = MarkdownToStructuredStep()
    ctx = make_ctx({"md_to_structured": {"rule_based": True}})

    result = step.run(ctx, [{"filename": "supply_chain_strategy.md", "markdown_text": md}])

    assert len(result) > 1
    assert result[0]["extraction_method"] == "rule_based"
    assert all(row.get("record_id") for row in result)
    assert {row["row_type"] for row in result} & {"section", "table_row", "rule"}
    assert result[0]["section_count"] > 0
    assert result[0]["doc_summary"]


def test_md_to_structured_rule_based_splits_markdown_table_rows():
    md = """# Suppliers

| supplier_id | supplier_name | score |
| --- | --- | --- |
| SUP001 | 天钢原材料有限公司 | 92 |
| SUP002 | 恒远包装材料公司 | 88 |
"""
    step = MarkdownToStructuredStep()
    ctx = make_ctx({"md_to_structured": {"rule_based": True}})

    result = step.run(ctx, [{"filename": "suppliers.md", "markdown_text": md}])

    table_rows = [row for row in result if row["row_type"] == "table_row"]
    assert len(table_rows) == 2
    assert table_rows[0]["supplier_id"] == "SUP001"
    assert table_rows[1]["supplier_name"] == "恒远包装材料公司"
