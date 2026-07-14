"""文档 → Markdown 转换 Step (策略模式)"""
from __future__ import annotations
import base64
import logging
import mimetypes
from pathlib import Path
from app.services.v2.pipeline.base import PipelineStep, PipelineContext

logger = logging.getLogger(__name__)


class DocumentToMarkdownStep(PipelineStep):
    """
    将文档转换为 Markdown 文本。

    spec 选项:
      strategy: "markitdown" | "ocr" | "vlm" (默认: "markitdown")
      model_id: str — VLM 策略使用的模型 ID (vlm 策略必填)

    input: data 的每个 row 形如 {"storage_uri": "s3://...", "filename": "..."}
    output: 每个 row 增加 "markdown_text" 字段
    """

    def run(self, ctx: PipelineContext, data: list[dict]) -> list[dict]:
        spec = ctx.spec.get("document_to_md", {})
        strategy = spec.get("strategy", "markitdown")

        result = []
        already_md = 0
        converted = 0

        for row in data:
            row = dict(row)

            # PRD media_reference: 添加来源文件引用字段
            if "filename" in row and "source_file" not in row:
                row["source_file"] = row["filename"]
            if ctx.dataset_id and "source_dataset_id" not in row:
                row["source_dataset_id"] = ctx.dataset_id

            # 已有非空 markdown_text 则跳过转换（来自 pipeline_run_task 预处理）
            if row.get("markdown_text"):
                row.setdefault("extraction_strategy", "passthrough")
                already_md += 1
                result.append(row)
                continue

            try:
                md = self._convert(row, strategy, spec, ctx)
                row["markdown_text"] = md
                row["extraction_strategy"] = strategy
                if ctx.meta.get("document_extraction_method"):
                    row["extraction_method"] = ctx.meta["document_extraction_method"]
                converted += 1
            except Exception as e:
                logger.warning(f"DocumentToMarkdown failed for {row.get('filename')}: {e}")
                row["markdown_text"] = ""
                row["extraction_error"] = str(e)
            result.append(row)

        ctx.meta["document_to_md"] = {
            "strategy": strategy,
            "processed": len(result),
            "converted": converted,
            "passthrough": already_md,
        }
        return result

    def _convert(self, row: dict, strategy: str, spec: dict, ctx: PipelineContext) -> str:
        filename = row.get("filename", "")
        content = row.get("content", b"")  # bytes 或字符串

        if strategy == "markitdown":
            return self._convert_markitdown(content, filename)
        elif strategy == "ocr":
            return self._convert_ocr(content, filename, spec, ctx)
        elif strategy == "vlm":
            return self._convert_vlm(content, filename, spec, ctx)
        else:
            raise ValueError(f"Unknown strategy: {strategy}")

    def _convert_markitdown(self, content: bytes | str, filename: str) -> str:
        """用 MarkItDown 将文档转换为 Markdown"""
        try:
            from markitdown import MarkItDown
            import tempfile
            import os
            md_converter = MarkItDown()

            if isinstance(content, bytes):
                # 创建临时文件
                suffix = Path(filename).suffix or ".bin"
                with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                    tmp.write(content)
                    tmp_path = tmp.name
                try:
                    result = md_converter.convert(tmp_path)
                    return result.text_content if hasattr(result, "text_content") else str(result)
                finally:
                    os.unlink(tmp_path)
            else:
                return str(content)
        except BaseException as e:
            logger.warning(f"MarkItDown conversion failed for {filename}: {e}")
            # MarkItDown 失败时回退为文本解码
            if isinstance(content, bytes):
                return content.decode("utf-8", errors="replace")
            return str(content)

    def _convert_ocr(self, content: bytes | str, filename: str, spec: dict | None = None,
                     ctx: PipelineContext | None = None) -> str:
        """Extract text from scanned images/PDFs with a configured OCR provider."""
        import os
        runtime = self._load_ocr_runtime_config((spec or {}).get("model_id"))
        provider = str((spec or {}).get("provider") or runtime.get("provider") or "easyocr").lower()
        enabled_env = os.getenv("ENABLE_OCR", "").lower() in ("1", "true", "yes")
        provider_env = os.getenv(f"ENABLE_{provider.upper()}", "").lower() in ("1", "true", "yes")
        if not enabled_env and not provider_env and not runtime.get("enabled"):
            logger.info("%s OCR disabled; enable it in OCR model config or set ENABLE_OCR=1", provider)
            if ctx:
                ctx.meta["document_extraction_method"] = f"{provider}_disabled"
            return ""

        if isinstance(content, str):
            return content

        image_bytes = content
        lower = filename.lower()
        if lower.endswith(".pdf") or content[:4] == b"%PDF":
            rendered, _ = self._render_pdf_first_page(content)
            image_bytes = rendered or content

        if provider == "paddleocr":
            return self._convert_paddleocr(image_bytes, filename, runtime, ctx)
        return self._convert_easyocr(image_bytes, filename, runtime, ctx)

    def _convert_easyocr(self, content: bytes, filename: str, runtime: dict,
                         ctx: PipelineContext | None = None) -> str:
        try:
            import easyocr
            import tempfile
            import os

            langs = runtime.get("langs") or runtime.get("lang") or ["ch_sim", "en"]
            if isinstance(langs, str):
                langs = [x.strip() for x in langs.replace(",", "\n").splitlines() if x.strip()]
            reader = easyocr.Reader(langs or ["ch_sim", "en"], gpu=runtime.get("device") == "gpu")
            suffix = Path(filename).suffix if Path(filename).suffix.lower() in {".png", ".jpg", ".jpeg", ".webp"} else ".png"
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            try:
                result = reader.readtext(tmp_path, detail=0, paragraph=True)
                if ctx:
                    ctx.meta["document_extraction_method"] = "easyocr"
                return "\n".join(str(item) for item in result if item)
            finally:
                os.unlink(tmp_path)
        except ImportError:
            logger.info("EasyOCR not available, returning empty string")
            if ctx:
                ctx.meta["document_extraction_method"] = "easyocr_unavailable"
            return ""
        except Exception as e:
            logger.warning(f"EasyOCR failed for {filename}: {e}")
            if ctx:
                ctx.meta["document_extraction_method"] = "easyocr_error"
            return ""

    def _convert_paddleocr(self, content: bytes, filename: str, runtime: dict,
                           ctx: PipelineContext | None = None) -> str:
        try:
            from paddleocr import PaddleOCR
            import tempfile
            import os

            ocr = PaddleOCR(use_angle_cls=True, lang=str(runtime.get("lang") or "ch"))
            suffix = Path(filename).suffix if Path(filename).suffix.lower() in {".png", ".jpg", ".jpeg", ".webp"} else ".png"
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            try:
                result = ocr.ocr(tmp_path, cls=True)
                lines = []
                for line in result:
                    if line:
                        for word_info in line:
                            if word_info and len(word_info) > 1:
                                lines.append(word_info[1][0])
                if ctx:
                    ctx.meta["document_extraction_method"] = "paddleocr"
                return "\n".join(lines)
            finally:
                os.unlink(tmp_path)
        except ImportError:
            logger.info("PaddleOCR not available, returning empty string")
            if ctx:
                ctx.meta["document_extraction_method"] = "paddleocr_unavailable"
            return ""
        except Exception as e:
            logger.warning(f"PaddleOCR failed for {filename}: {e}")
            if ctx:
                ctx.meta["document_extraction_method"] = "paddleocr_error"
            return ""

    def _load_ocr_runtime_config(self, model_id: str | None = None) -> dict:
        try:
            from app.database import SessionLocal
            from app.models.model_config import ModelConfig

            db = SessionLocal()
            try:
                query = db.query(ModelConfig).filter(ModelConfig.config_type == "ocr")
                cfg = query.filter(ModelConfig.id == model_id).first() if model_id else None
                if not cfg:
                    cfg = query.filter(ModelConfig.provider == "easyocr").order_by(ModelConfig.updated_at.desc()).first()
                if not cfg:
                    cfg = query.order_by(ModelConfig.updated_at.desc()).first()
                if not cfg:
                    return {}
                options = dict(cfg.options or {})
                options["provider"] = cfg.provider
                return options
            finally:
                db.close()
        except Exception:
            return {}

    def _convert_vlm(self, content: bytes | str, filename: str, spec: dict, ctx: PipelineContext) -> str:
        """Use a configured vision-capable LLM to convert image/PDF content to Markdown."""
        if isinstance(content, str):
            return content

        image_bytes, mime_type = self._vlm_input_image(content, filename)
        if not image_bytes:
            structured = self._extract_office_markdown(content, filename)
            if structured:
                ctx.meta["document_extraction_method"] = "office_structured"
                return structured
            logger.info("VLM input is not directly renderable; falling back to MarkItDown for %s", filename)
            ctx.meta["document_extraction_method"] = "markitdown_fallback"
            return self._convert_markitdown(content, filename)

        model_config = self._load_vlm_model_config(spec.get("model_id"))
        if not model_config:
            logger.info("No VLM model config found; falling back to MarkItDown for %s", filename)
            return self._convert_markitdown(content, filename)

        try:
            ctx.meta["document_extraction_method"] = "vlm"
            return self._call_vlm(model_config, image_bytes, mime_type, filename, spec)
        except Exception as e:
            logger.warning("VLM extraction failed for %s: %s", filename, e)
            ctx.meta["document_extraction_method"] = "markitdown_fallback"
            return self._convert_markitdown(content, filename)

    def _vlm_input_image(self, content: bytes, filename: str) -> tuple[bytes | None, str | None]:
        mime_type = mimetypes.guess_type(filename)[0] or ""
        lower = filename.lower()
        if mime_type.startswith("image/") or lower.endswith((".png", ".jpg", ".jpeg", ".webp")):
            return content, mime_type or "image/png"
        if lower.endswith(".pdf") or content[:4] == b"%PDF":
            return self._render_pdf_first_page(content)
        return None, None

    def _extract_office_markdown(self, content: bytes, filename: str) -> str:
        lower = filename.lower()
        if lower.endswith(".pptx"):
            return self._extract_pptx_markdown(content)
        if lower.endswith(".docx"):
            return self._extract_docx_markdown(content)
        return ""

    def _extract_pptx_markdown(self, content: bytes) -> str:
        try:
            import io
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            lines: list[str] = []
            for slide_idx, slide in enumerate(prs.slides, start=1):
                lines.append(f"# Slide {slide_idx}")
                text_blocks: list[str] = []
                table_blocks: list[str] = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        text = "\n".join(
                            p.text.strip()
                            for p in shape.text_frame.paragraphs
                            if p.text and p.text.strip()
                        )
                        if text:
                            text_blocks.append(text)
                    if getattr(shape, "has_table", False):
                        table_blocks.append(self._pptx_table_to_markdown(shape.table))
                lines.extend(text_blocks)
                lines.extend(t for t in table_blocks if t)
                lines.append("")
            return "\n".join(lines).strip()
        except Exception as e:
            logger.info("PPTX structured extraction failed: %s", e)
            return ""

    def _pptx_table_to_markdown(self, table) -> str:
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip().replace("\n", " ") for cell in row.cells])
        return self._rows_to_markdown_table(rows)

    def _extract_docx_markdown(self, content: bytes) -> str:
        try:
            import io
            from docx import Document

            doc = Document(io.BytesIO(content))
            lines: list[str] = []
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                style_name = (paragraph.style.name or "").lower() if paragraph.style else ""
                if "heading" in style_name:
                    level = "".join(ch for ch in style_name if ch.isdigit()) or "1"
                    lines.append(f"{'#' * max(1, min(int(level), 6))} {text}")
                else:
                    lines.append(text)
            for table in doc.tables:
                rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells] for row in table.rows]
                table_md = self._rows_to_markdown_table(rows)
                if table_md:
                    lines.append(table_md)
            return "\n\n".join(lines).strip()
        except Exception as e:
            logger.info("DOCX structured extraction failed: %s", e)
            return ""

    def _rows_to_markdown_table(self, rows: list[list[str]]) -> str:
        if not rows:
            return ""
        width = max(len(row) for row in rows)
        normalized = [row + [""] * (width - len(row)) for row in rows]
        header = normalized[0]
        separator = ["---"] * width
        body = normalized[1:]

        def fmt(row: list[str]) -> str:
            return "| " + " | ".join(cell or "" for cell in row) + " |"

        return "\n".join([fmt(header), fmt(separator), *(fmt(row) for row in body)])

    def _render_pdf_first_page(self, content: bytes) -> tuple[bytes | None, str | None]:
        try:
            import fitz

            doc = fitz.open(stream=content, filetype="pdf")
            try:
                if doc.page_count == 0:
                    return None, None
                page = doc.load_page(0)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                return pix.tobytes("png"), "image/png"
            finally:
                doc.close()
        except Exception as e:
            logger.info("PDF render for VLM failed: %s", e)
            return None, None

    def _load_vlm_model_config(self, model_id: str | None = None):
        try:
            from app.services.model_config_selector import select_llm_model_config
            return select_llm_model_config(
                model_id=model_id,
                purpose_tags=("VLM提取", "vlm", "vision"),
                allow_vlm=True,
            )
        except Exception:
            return None

    def _call_vlm(self, model_config, image_bytes: bytes, mime_type: str | None, filename: str, spec: dict) -> str:
        from app.services.model_config_selector import llm_call_kwargs

        call_kwargs = llm_call_kwargs(model_config)
        if not call_kwargs:
            raise RuntimeError("No model configured for VLM extraction")
        prompt = spec.get("prompt") or (
            "Extract all visible text, tables, headings, entities, dates, numbers and business rules from this document image. "
            "Return concise Markdown only. Preserve tables as Markdown tables when possible."
        )
        encoded = base64.b64encode(image_bytes).decode("ascii")
        if model_config.provider == "anthropic":
            api_key = call_kwargs["api_key"]
            model_name = call_kwargs["model"]
            import anthropic

            client = anthropic.Anthropic(api_key=api_key)
            response = client.messages.create(
                model=model_name,
                max_tokens=int((model_config.options or {}).get("max_tokens", 4096)),
                system="You are a document vision extraction engine. Return Markdown only.",
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image", "source": {
                            "type": "base64",
                            "media_type": mime_type or "image/png",
                            "data": encoded,
                        }},
                        {"type": "text", "text": f"File: {filename}\n{prompt}"},
                    ],
                }],
            )
            return response.content[0].text if response.content else ""

        from app.services.llm_service import _call_llm

        data_url = f"data:{mime_type or 'image/png'};base64,{encoded}"
        messages = [
            {"role": "system", "content": "You are a document vision extraction engine. Return Markdown only."},
            {"role": "user", "content": [
                {"type": "text", "text": f"File: {filename}\n{prompt}"},
                {"type": "image_url", "image_url": {"url": data_url}},
            ]},
        ]
        return _call_llm(
            call_kwargs["provider"],
            call_kwargs["api_key"],
            call_kwargs["api_base"],
            call_kwargs["model"],
            messages,
            json_mode=False,
        )
